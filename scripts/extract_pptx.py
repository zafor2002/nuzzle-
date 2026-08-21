import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
import os

pptx_path = r'e:\nuzzle project\forntenduse vue js\presentation\Nuzzle_UIUX_Presentation (3) (2).pptx'
prs = Presentation(pptx_path)

print(f'Total slides: {len(prs.slides)}')
for i, slide in enumerate(prs.slides):
    print(f'\n=== SLIDE {i+1} ===')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    # print font info for first run if available
                    color_info = ''
                    size_info = ''
                    bold_info = ''
                    if para.runs:
                        r = para.runs[0]
                        if r.font.color and r.font.color.type is not None:
                            try:
                                color_info = f' [color=#{r.font.color.rgb}]'
                            except:
                                pass
                        if r.font.size:
                            size_info = f' [size={r.font.size.pt}pt]'
                        if r.font.bold:
                            bold_info = ' [BOLD]'
                    print(f'  TEXT: {text}{color_info}{size_info}{bold_info}')
        # Check fill colors of shapes
        if shape.shape_type in [1, 6]:  # rectangle-like shapes
            try:
                fill = shape.fill
                if fill.type is not None and fill.fore_color.type is not None:
                    try:
                        print(f'  SHAPE fill=#{fill.fore_color.rgb} w={round(shape.width/914400,2)}" h={round(shape.height/914400,2)}"')
                    except:
                        pass
            except:
                pass

print('\n\nDone.')
